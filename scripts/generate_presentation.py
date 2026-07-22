import sys
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors matching SecurePrompt UI
    bg_color = RGBColor(11, 15, 25)        # #0B0F19
    white = RGBColor(255, 255, 255)       # #FFFFFF
    indigo = RGBColor(129, 140, 248)      # #818CF8
    emerald = RGBColor(52, 211, 153)      # #34D399
    slate_text = RGBColor(148, 163, 184)  # #94A3B8
    card_bg = RGBColor(22, 28, 45)         # #161C2D
    
    blank_layout = prs.slide_layouts[6]
    
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
    def add_header(slide, title, category="SECUREPROMPT AI GATEWAY"):
        # Category label
        add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3), 
                    category, font_size=10, font_color=indigo, bold=True)
        # Main slide title
        add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8), 
                    title, font_size=28, font_color=white, bold=True)
        
    def add_textbox(slide, left, top, width, height, text, font_size, font_color, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Calibri"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = align
        return txBox

    def add_bullet_list(slide, left, top, width, height, items, font_size=15):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "•  " + item
            p.font.name = "Calibri"
            p.font.size = Pt(font_size)
            p.font.color.rgb = slate_text
            p.space_after = Pt(10)
            
    def insert_image_safe(slide, filename, left, top, width, height):
        screenshot_path = Path("docs/screenshots") / filename
        if screenshot_path.exists():
            slide.shapes.add_picture(str(screenshot_path), left, top, width=width, height=height)
            # Add subtle border to screenshot
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            rect.fill.background()
            rect.line.color.rgb = RGBColor(51, 65, 85)
            rect.line.width = Pt(1.5)
            print(f"Embedded: {filename}")
        else:
            # Placeholder if screenshot is missing
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = card_bg
            rect.line.color.rgb = RGBColor(51, 65, 85)
            add_textbox(slide, left + Inches(0.2), top + (height/2) - Inches(0.3), width - Inches(0.4), Inches(0.6),
                        f"Screenshot: {filename}\n(Not found in docs/screenshots)", font_size=12, font_color=slate_text, align=PP_ALIGN.CENTER)

    # ----------------------------------------------------
    # SLIDE 1: TITLE SLIDE
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)
    
    # Title Box in center
    add_textbox(slide1, Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.5), 
                "SecurePrompt AI Gateway", font_size=52, font_color=white, bold=True)
    
    add_textbox(slide1, Inches(1.0), Inches(3.6), Inches(11.3), Inches(0.8), 
                "Enterprise Zero-Trust AI DLP Firewall & Role-Based Access Control (RBAC) Gateway", 
                font_size=20, font_color=indigo, bold=False)
    
    # Divider line
    divider = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.5), Inches(4.0), Inches(0.04))
    divider.fill.solid()
    divider.fill.fore_color.rgb = emerald
    divider.line.fill.background()
    
    add_textbox(slide1, Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.5), 
                "Designed and Developed by: Prathyusha (IAM & AI Security Specialist)\n"
                "Showcasing skills in: Identity & Access Management, Data Loss Prevention, & GenAI Gateways", 
                font_size=14, font_color=slate_text)

    # ----------------------------------------------------
    # SLIDE 2: THE BUSINESS CHALLENGE (GenAI Data Leakage)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2)
    add_header(slide2, "The Enterprise AI Security Challenge")
    
    bullets = [
        "Uncontrolled LLM Adoption: Employees and developers paste proprietary code, cloud credentials, and corporate reports into ChatGPT/Claude.",
        "Corporate Data Leakage (DLP): Third-party AI providers retain prompts for training, creating compliance violations (GDPR, DPDP, HIPAA, PCI-DSS).",
        "Credential Exposure: API keys, cloud connections, and GitHub/Slack tokens pasted into prompts expose internal infrastructure.",
        "Lack of Access Control (IAM): Traditional firewalls cannot read or govern prompt payload contexts, leaving security teams blind.",
        "The Solution: SecurePrompt AI Gateway sits between the user and the LLM, analyzing every request in real time BEFORE it leaves the company boundary."
    ]
    add_bullet_list(slide2, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), bullets, font_size=16)

    # ----------------------------------------------------
    # SLIDE 3: SYSTEM ARCHITECTURE & DATA FLOW
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3)
    add_header(slide3, "System Architecture & End-to-End Data Flow")
    
    # Left Column: Features
    bullets_arch = [
        "Frontend Client: Single-page React application compiled via Babel, styled with Tailwind CSS, utilizing Chart.js for visualization.",
        "FastAPI Firewall Gateway: High-performance, asynchronous REST backend responsible for inspection, policy evaluation, and routing.",
        "Multi-Stage Scanner Pipeline: Prompts and attachments undergo sequential checks (PII, secrets, jailbreaks, prompt injection).",
        "Remediation & Hashing: High-risk entries are blocked; remediable items are masked/redacted automatically using Scrypt entropy-safe hashes.",
        "Observability DB: Audit logs and policy violations logged to a central PostgreSQL (Supabase) database with transactional connection pooling."
    ]
    add_bullet_list(slide3, Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0), bullets_arch, font_size=14)
    
    # Right Column: Architecture Diagram
    add_textbox(slide3, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.4), "System Topology Flow", font_size=14, font_color=indigo, bold=True)
    
    # We can place shapes to draw a beautiful block architecture flow!
    flow_steps = [
        ("User Prompt / File Upload", Inches(2.3)),
        ("FastAPI /analyze API", Inches(3.2)),
        ("DLP Pipeline (PII/Secrets/Jailbreak Scans)", Inches(4.1)),
        ("Policy & Risk Engine (policies.yaml)", Inches(5.0)),
        ("Dynamic Outbound LLM Router (Groq/OpenAI)", Inches(5.9)),
    ]
    for text, top_pos in flow_steps:
        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(top_pos), Inches(5.5), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = card_bg
        box.line.color.rgb = RGBColor(51, 65, 85)
        # Add text
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Calibri"
        p.font.size = Pt(13)
        p.font.color.rgb = white
        p.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 4: THE ZERO-TRUST SCANNER PIPELINE (Inspection Logic)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4)
    add_header(slide4, "Zero-Trust Scanner Pipeline & Pattern Logic")
    
    bullets_scanner = [
        "AWS & Cloud Keys: Matches AWS_ACCESS_KEY (using AKIA prefix checking) and AWS_SECRET_KEY pattern syntax.",
        "GitHub & Slack Tokens: Scans for oauth, personal access tokens (ghp_, ghx_), and Slack bots/app tokens.",
        "PII & PHI Detections: Captures regional identifiers (PAN Card, Aadhaar Card, SSN) and medical reports (PHI).",
        "OpenAI API Keys (Modern sk-proj-): Enhanced regex pattern 'sk-[A-Za-z0-9_\\-]{20,}' successfully detects modern OpenAI keys containing dashes.",
        "Adversarial Protection: Detects prompt injection phrases and adversarial jailbreaks designed to bypass safety filters."
    ]
    add_bullet_list(slide4, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), bullets_scanner, font_size=14)
    
    # Right column: code block placeholder or mock detection screenshot
    insert_image_safe(slide4, "media__1784444045087.png", Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.5))

    # ----------------------------------------------------
    # SLIDE 5: ENTERPRISE DLP POLICY ENGINE
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5)
    add_header(slide5, "Enterprise-Grade DLP Policy Configurations")
    
    bullets_policies = [
        "Declarative Control (policies.yaml): Policies are configured in YAML and reloaded on the fly without code changes or restarts.",
        "BLOCK_PROPRIETARY_CODE: Detects programming syntax/keywords or file extensions (.py, .js, .java) and restricts code upload.",
        "BLOCK_CLOUD_SECRETS: Catches cloud provider keys, API tokens, and database connection strings.",
        "BLOCK_HR_PAYROLL: Restricts salary spreadsheets, compensation letters, and payroll databases.",
        "BLOCK_INFRA_ARCHITECTURE: Blocks upload of system architectures, network topology documents, or database design files."
    ]
    add_bullet_list(slide5, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), bullets_policies, font_size=15)

    # ----------------------------------------------------
    # SLIDE 6: SECURE SIGN-UP & SIGN-IN FLOW (IAM)
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6)
    add_header(slide6, "IAM: Secure Registration & Authentication")
    
    bullets_auth = [
        "Enterprise Portal: Sleek Sign In / Create Account workflow secures the gateway from unauthorized public access.",
        "Native Hashing (scrypt): Passwords are fully hashed with a randomized salt using the built-in hashlib.scrypt (memory-hard, recommended by OWASP).",
        "Zero-Dependency Security: Solves compatibility issues and security vulnerabilities of legacy packages by using Python native libraries.",
        "Enterprise Role Selection: Users register with distinct roles - 'User (Developer)' or 'Security Admin' - which dictate dashboard privileges."
    ]
    add_bullet_list(slide6, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), bullets_auth, font_size=14)
    
    # Embed login/signup screenshot (media__1784445820655.png or media__1784446682695.png)
    insert_image_safe(slide6, "media__1784446682695.png", Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.5))

    # ----------------------------------------------------
    # SLIDE 7: ROLE-BASED ACCESS CONTROL (RBAC)
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7)
    add_header(slide7, "IAM: Role-Based Observability & Auditing")
    
    bullets_rbac = [
        "Greeting & Role Badge: Greeting ('Hello prathyusha') and active role displayed in the top right ('developer' or 'security_admin').",
        "Developer Role (Least Privilege): Data is strictly partitioned. Developers can only view their own dashboard, audit logs, and reports.",
        "Security Admin (Auditor): Bypasses filters to see the aggregated logs of all developers, maintaining full oversight.",
        "Safe Gateway Remediation: If a prompt violates policy, developers can click 'Mask Sensitive Data (Auto Fix)' to send a safe, redacted version."
    ]
    add_bullet_list(slide7, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), bullets_rbac, font_size=14)
    
    # Embed gateway with user/role badge screenshot (media__1784444786623.png)
    insert_image_safe(slide7, "media__1784444786623.png", Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.5))

    # ----------------------------------------------------
    # SLIDE 8: DYNAMIC EXPLORATORY DATA ANALYSIS (EDA)
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_background(slide8)
    add_header(slide8, "Real-Time EDA Analysis & Observability")
    
    bullets_eda = [
        "Date-Filtered Analysis: Dynamic date pickers calculate analytics reports for specific compliance windows.",
        "Line Chart: Illustrates daily request volume and average risk trends.",
        "Doughnut Chart: Shows the decision distribution (Allowed, Warned, Blocked).",
        "Bar Chart: Renders entity detections overview (Aadhaar, PAN, SSN, Secrets, Code).",
        "Automated Anomalies & Insights: Highlights outliers (e.g. 'Allowed prompts are on average 46% longer than blocked ones').",
        "Supabase Connection Pooling: Transaction pooling handles high concurrent DB queries during analytics compilation."
    ]
    add_bullet_list(slide8, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), bullets_eda, font_size=14)
    
    # Embed analysis tab charts screenshot (media__1784442779392.png)
    insert_image_safe(slide8, "media__1784442779392.png", Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.5))

    # ----------------------------------------------------
    # SLIDE 9: SUMMARY & SECURE AI ROADMAP
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_background(slide9)
    add_header(slide9, "Summary & Career Roadmap: Security with AI")
    
    bullets_roadmap = [
        "Full-Stack Foundation: Built a high-performance, containerized, and deployed firewall bridging Web App engineering and AI governance.",
        "Core Skill Demonstration: Showcases strong proficiency in IAM (RBAC, Scrypt Hashing, Token Contexts) and DLP (Regex, Heuristics, Policy Engine).",
        "Future AI Security Use Cases to Propose to Manager:",
        "   1. LLM Evaluation & Guardrails: Add alignment scanners (hallucination checks, toxic output filtering, prompt alignment).",
        "   2. Vector DB Sanitization: Implement safe embeddings (strip PII before converting prompts to vectors for retrieval-augmented generation).",
        "   3. API Key Vault Integration: Store target LLM API keys in HashiCorp Vault or AWS Secrets Manager, injected dynamically by the gateway."
    ]
    add_bullet_list(slide9, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), bullets_roadmap, font_size=15)

    # Save presentation
    output_path = Path("docs/SecurePrompt_AI_Gateway_Presentation.pptx")
    prs.save(str(output_path))
    print(f"Presentation saved successfully at: {output_path.absolute()}")

if __name__ == "__main__":
    create_presentation()
